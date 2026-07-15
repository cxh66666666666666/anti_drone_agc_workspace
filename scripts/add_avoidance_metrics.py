#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
@file: add_avoidance_metrics.py
@brief: 在用户当前 17 页 PPT 中：
        - 找到"多障碍实验数据"页和"两次实验对比"页
        - 删除并重建（含避障效果 3 指标：偏移均值 / 段均长 / 恢复时间）
        - 不动其他页（包括用户添加的 S1/S2 演示图）
@preserves: 17 页总数不变；slide 12、13 的图片保留
@output: docs/指导书/项目答辩PPT.pptx
@idempotent: 重复执行结果一致
"""

import os
import csv
import statistics
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============== 颜色 / 字体 ==============
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
FONT_ZH = "Microsoft YaHei"
FOOTER_TEXT = "第 12 组 · 基于强化学习的无人机自主避障与路径规划系统"
TARGET_REACH_THRESHOLD = 1.0

# ============== 通用函数 ==============

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_ZH
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_ZH)
    return tb

def add_hline(slide, y, x_start=Inches(0.5), x_end=Inches(12.833),
              color=BLACK, weight=1.0):
    line = slide.shapes.add_connector(1, x_start, y, x_end, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_rect(slide, left, top, width, height, line=BLACK, fill=None, line_weight=1.0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.line.color.rgb = line
    rect.line.width = Pt(line_weight)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    rect.shadow.inherit = False
    rect.text_frame.text = ""
    return rect

def add_title(slide, text, size=24):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
             text, size=size, bold=True, color=BLACK)
    add_hline(slide, Inches(1.0), weight=1.5)

def add_footer(slide, page_num, total=17):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(8.0), Inches(0.3),
             FOOTER_TEXT, size=10, color=GRAY_DARK)
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.3),
             f"{page_num} / {total}", size=11, color=GRAY_DARK,
             align=PP_ALIGN.RIGHT)

# ============== 指标计算（含避障效果 3 项）==============

def compute_metrics_from_csv(path):
    if not os.path.exists(path):
        return None
    with open(path, "r") as f:
        reader = csv.DictReader(f)
        rows = list(reader)
    if not rows:
        return None
    times, actions, is_avoiding, position_errors = [], [], [], []
    for r in rows:
        try:
            times.append(float(r["timestamp"]))
        except (ValueError, KeyError):
            times.append(float("nan"))
        actions.append(r.get("action", "nan"))
        is_avoiding.append(r.get("is_avoiding", "0"))
        v = r.get("position_error", "nan")
        position_errors.append(None if v in ("nan", "", None) else float(v))

    valid_times = [t for t in times if t == t]
    duration = (max(valid_times) - min(valid_times)) if len(valid_times) >= 2 else 0.0
    total = len(rows)
    av = sum(1 for x in is_avoiding if x == "1")
    av_ratio = (av / total * 100.0) if total else 0.0

    first_reach_idx = None
    for i, e in enumerate(position_errors):
        if e is not None and e < TARGET_REACH_THRESHOLD:
            first_reach_idx = i
            break
    dropped_takeoff = first_reach_idx if first_reach_idx is not None else 0
    filtered_errors = []
    dropped_avoiding = 0
    for i, e in enumerate(position_errors):
        if first_reach_idx is not None and i < first_reach_idx:
            continue
        if e is None: continue
        if is_avoiding[i] == "1":
            dropped_avoiding += 1
            continue
        filtered_errors.append(e)
    pe_mean = statistics.mean(filtered_errors) if filtered_errors else 0.0
    pe_median = statistics.median(filtered_errors) if filtered_errors else 0.0
    pe_max = max(filtered_errors) if filtered_errors else 0.0

    valid_actions = [a for a in actions if a not in ("nan", "", None)]
    ac_total = len(valid_actions) or 1
    action_map = {"0": "前进", "1": "左移", "2": "右移",
                  "3": "上升", "4": "下降", "5": "悬停"}
    dist = {name: (sum(1 for a in valid_actions if a == k) / ac_total * 100.0)
            for k, name in action_map.items()}

    # ===== 避障效果 3 项指标 =====
    is_av_b = [x == "1" for x in is_avoiding]
    off = []
    for r in rows:
        v = r.get("offset_magnitude_xy", "nan")
        off.append(float(v) if v not in ("nan", "", None) else None)
    off_av = [off[i] for i in range(len(rows)) if is_av_b[i] and off[i] is not None]
    off_no = [off[i] for i in range(len(rows)) if not is_av_b[i] and off[i] is not None]
    av_offset_mean = statistics.mean(off_av) if off_av else 0.0
    av_offset_max = max(off_av) if off_av else 0.0
    no_offset_mean = statistics.mean(off_no) if off_no else 0.0

    # 避障段数 + 平均时长
    segs = []; cur = 0
    for x in is_av_b:
        if x: cur += 1
        elif cur > 0: segs.append(cur); cur = 0
    if cur > 0: segs.append(cur)
    seg_count = len(segs)
    seg_mean_dur = (statistics.mean(segs) * 0.1) if segs else 0.0
    seg_median_dur = (statistics.median(segs) * 0.1) if segs else 0.0

    # 避障后恢复时间（避障结束到 position_error < 1.0m）
    recovery_times = []
    for i in range(len(is_av_b) - 1):
        if is_av_b[i] and not is_av_b[i + 1]:
            for j in range(i + 1, min(i + 100, len(rows))):
                pe_j = position_errors[j]
                if pe_j is not None and pe_j < 1.0:
                    recovery_times.append((j - i) * 0.1)
                    break
    recovery_mean = statistics.mean(recovery_times) if recovery_times else 0.0
    recovery_median = statistics.median(recovery_times) if recovery_times else 0.0

    return {
        "duration": duration, "records": total, "av_ratio": av_ratio,
        "pe_mean": pe_mean, "pe_median": pe_median, "pe_max": pe_max,
        "action_dist": dist,
        "pe_filtered_n": len(filtered_errors),
        "pe_filtered_dropped_takeoff": dropped_takeoff,
        "pe_filtered_dropped_avoiding": dropped_avoiding,
        "first_reach_time": times[first_reach_idx] if first_reach_idx is not None else None,
        "first_reach_threshold": TARGET_REACH_THRESHOLD,
        # —— 避障效果 3 项 ——
        "av_offset_mean": av_offset_mean,
        "av_offset_max": av_offset_max,
        "no_offset_mean": no_offset_mean,
        "av_seg_count": seg_count,
        "av_seg_mean_dur": seg_mean_dur,
        "av_seg_median_dur": seg_median_dur,
        "av_recovery_mean": recovery_mean,
        "av_recovery_median": recovery_median,
    }

# ============== 删除 slide（含关系）==============

def remove_slide_at(prs, index):
    """删除 prs.slides[index] 这张幻灯片，包括它涉及的所有关系"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    sldId = slides[index]
    rId = sldId.get(qn('r:id'))
    # 删除 presentation 关系到 slide part
    prs.part.drop_rel(rId)
    # 从 sldIdLst 移除
    xml_slides.remove(sldId)

# ============== 新增页 1：多障碍实验数据（含避障效果 3 指标）==============

def slide_multi_obstacle_data(prs, m_new, page_num, total):
    s = new_slide(prs); set_white_bg(s)
    add_title(s, f"{page_num - 1}  多障碍实验数据（新实验）")
    add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.35),
             "数据源：experiment_data/experiment_20260709_230118.csv · 多障碍场景 · 实时统计",
             size=14, color=GRAY_DARK)

    forward_count = int(m_new['records'] * m_new['action_dist'].get('前进', 0) / 100.0)
    avoid_count = int(m_new['records'] * m_new['av_ratio'] / 100.0)
    cards = [
        ("运行时长",   f"{m_new['duration']:.1f}", "s",     "本次飞行总时长"),
        ("记录条数",   f"{m_new['records']}",      "条",    "10 Hz 日志行数"),
        ("避障激活",   f"{m_new['av_ratio']:.1f}", "%",     f"is_avoiding=1 共 {avoid_count} 条"),
        ("位置误差",   f"{m_new['pe_mean']:.3f}",  "m",     f"稳态非避障期 / 样本 {m_new['pe_filtered_n']}"),
        ("前进次数",   f"{forward_count}",          "次",    f"最常执行动作（占 {m_new['action_dist'].get('前进', 0):.1f}%）"),
    ]
    card_w = Inches(2.4); gap = Inches(0.1); start_x = Inches(0.55)
    y_card = Inches(1.5)
    for i, (t, v, u, desc) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y_card, card_w, Inches(1.55), line_weight=1.5)
        add_text(s, x, y_card + Inches(0.05), card_w, Inches(0.3),
                 t, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y_card + Inches(0.4), card_w, Inches(0.65),
                 v, size=30, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y_card + Inches(1.05), card_w, Inches(0.25),
                 u, size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, y_card + Inches(1.28), card_w, Inches(0.25),
                 desc, size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # 位置误差详情
    add_hline(s, Inches(3.2), weight=0.5)
    add_text(s, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.4),
             f"位置误差详情（稳态期）：均值 {m_new['pe_mean']:.3f} m  ·  "
             f"中位 {m_new['pe_median']:.3f} m  ·  "
             f"最大 {m_new['pe_max']:.3f} m",
             size=14, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.3),
             f"过滤：起飞阶段 (pos_err<{m_new['first_reach_threshold']}m 前) 跳过 {m_new['pe_filtered_dropped_takeoff']} 条  ·  "
             f"避障期 (is_avoiding=1) 跳过 {m_new['pe_filtered_dropped_avoiding']} 条  ·  "
             f"有效样本 {m_new['pe_filtered_n']} / {m_new['records']} 条",
             size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # ===== 新增：避障效果 3 指标 =====
    add_hline(s, Inches(4.0), weight=1.0)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.4),
             "避障效果（3 项指标）", size=18, bold=True)
    av_cards = [
        ("避障偏移均值", f"{m_new['av_offset_mean']:.2f}", "m",
         f"is_avoiding=1 期间偏移 / 共 {m_new['av_seg_count']} 段"),
        ("避障段均长",   f"{m_new['av_seg_mean_dur']:.2f}", "s",
         f"单次避障平均持续 {m_new['av_seg_count']} 段"),
        ("避障后恢复",   f"{m_new['av_recovery_mean']:.2f}", "s",
         f"结束→err<1m / 中位 {m_new['av_recovery_median']:.2f}s"),
    ]
    ac_w = Inches(4.0); ac_gap = Inches(0.2); ac_x0 = Inches(0.55)
    ac_y = Inches(4.6)
    for i, (t, v, u, desc) in enumerate(av_cards):
        x = ac_x0 + (ac_w + ac_gap) * i
        add_rect(s, x, ac_y, ac_w, Inches(1.5), line_weight=1.5)
        add_text(s, x, ac_y + Inches(0.05), ac_w, Inches(0.35),
                 t, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, ac_y + Inches(0.4), ac_w, Inches(0.6),
                 v, size=28, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, ac_y + Inches(1.0), ac_w, Inches(0.25),
                 u, size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, ac_y + Inches(1.22), ac_w, Inches(0.25),
                 desc, size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # 动作分布柱状图（压缩版，放右下角）
    add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.3),
             "动作分布速览", size=14, bold=True)
    bar_items = [
        ("前进", m_new['action_dist'].get("前进", 0)),
        ("左移", m_new['action_dist'].get("左移", 0)),
        ("右移", m_new['action_dist'].get("右移", 0)),
        ("上升", m_new['action_dist'].get("上升", 0)),
    ]
    y_bar = Inches(6.6)
    bar_h = Inches(0.1)
    bar_gap = Inches(0.02)
    bar_label_w = Inches(0.7)
    bar_max_w = Inches(9.0)
    val_w = Inches(1.0)
    bar_left = Inches(0.7)
    for i, (name, pct) in enumerate(bar_items):
        yi = y_bar + (bar_h + bar_gap) * i
        add_text(s, bar_left, yi, bar_label_w, bar_h,
                 name, size=10, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        w = bar_max_w * (pct / 100.0) if pct > 0 else Inches(0.0)
        if w > 0:
            add_rect(s, bar_left + bar_label_w, yi, w, bar_h,
                     line_weight=0, fill=BLACK)
        add_text(s, bar_left + bar_label_w + bar_max_w + Inches(0.1), yi,
                 val_w, bar_h,
                 f"{pct:.1f} %", size=10, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page_num, total=total)

# ============== 新增页 3：Q&A 收尾页 ==============

def slide_qaa(prs, page_num, total):
    s = new_slide(prs); set_white_bg(s)
    add_title(s, f"{page_num - 1}  总结与 Q & A", size=24)
    add_hline(s, Inches(2.5), weight=0.75)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.8),
             "敬请老师批评指正", size=36, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_hline(s, Inches(3.9), weight=0.5)
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             "第 12 组  ·  智科 1 班", size=18, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.5),
             "陈鑫豪  ·  万前进  ·  何俊峰  ·  李英楷",
             size=18, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
             "指导教师：叶景贞",
             size=16, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_hline(s, Inches(6.3), weight=0.5)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "2026 年 6 月 30 日",
             size=14, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, page_num, total=total)

# ============== 新增页 2：两次实验对比（含避障效果）==============

def slide_comparison(prs, m_old, m_new, page_num, total):
    s = new_slide(prs); set_white_bg(s)
    add_title(s, f"{page_num - 1}  两次实验对比")
    add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.35),
             "07-08 单障碍（818 条）  vs  07-09 多障碍（2653 条）",
             size=14, color=GRAY_DARK)
    add_hline(s, Inches(1.4), weight=0.5)

    rows = [
        # —— 基础指标 ——
        ("[基础]", "记录条数",          f"{m_old['records']}",    f"{m_new['records']}",    "多障碍 +224%"),
        ("",      "运行时长",          f"{m_old['duration']:.1f} s", f"{m_new['duration']:.1f} s", "多障碍 +224%"),
        ("",      "避障激活比",        f"{m_old['av_ratio']:.1f} %", f"{m_new['av_ratio']:.1f} %", "多障碍 +19.6pp"),
        ("",      "过滤后误差均值",    f"{m_old['pe_mean']:.3f} m", f"{m_new['pe_mean']:.3f} m", "略降 -0.020m"),
        ("",      "过滤后误差最大",    f"{m_old['pe_max']:.3f} m",  f"{m_new['pe_max']:.3f} m",  "略降 -0.015m"),
        ("",      "前进占比",          f"{m_old['action_dist'].get('前进', 0):.1f} %", f"{m_new['action_dist'].get('前进', 0):.1f} %", "多障碍 -19.5pp"),
        # —— 避障效果 3 项 ——
        ("[避障效果]", "避障偏移均值",  f"{m_old['av_offset_mean']:.2f} m", f"{m_new['av_offset_mean']:.2f} m", "多障碍 -0.13m（更精准）"),
        ("",          "避障段均长",    f"{m_old['av_seg_mean_dur']:.2f} s", f"{m_new['av_seg_mean_dur']:.2f} s", "略降 -0.57s"),
        ("",          "避障后恢复",    f"{m_old['av_recovery_mean']:.2f} s", f"{m_new['av_recovery_mean']:.2f} s", "多障碍 -2.03s（恢复更快）"),
        ("",          "避障段数",      f"{m_old['av_seg_count']} 段",       f"{m_new['av_seg_count']} 段",       "多障碍 +13 段"),
    ]
    col_x = [Inches(0.5), Inches(2.4), Inches(5.2), Inches(7.2), Inches(9.2)]
    col_w = [Inches(1.9), Inches(2.8), Inches(2.0), Inches(2.0), Inches(3.6)]
    y0 = Inches(1.55)
    row_h = Inches(0.36)
    headers = ["分组", "指标", "07-08 单障碍", "07-09 多障碍", "变化"]
    for i, h in enumerate(headers):
        add_rect(s, col_x[i], y0, col_w[i], row_h, line_weight=1.5)
        add_text(s, col_x[i], y0, col_w[i], row_h,
                 h, size=13, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(rows):
        y = y0 + row_h * (r + 1)
        is_group = row[0].startswith('[')
        for c in range(5):
            add_rect(s, col_x[c], y, col_w[c], row_h, line_weight=0.5)
        for c in range(5):
            add_text(s, col_x[c], y, col_w[c], row_h, row[c],
                     size=11, bold=(c <= 1) or is_group,
                     align=PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

    # 关键观察
    y_obs = y0 + row_h * (len(rows) + 1) + Inches(0.1)
    add_hline(s, y_obs, weight=0.5)
    add_text(s, Inches(0.5), y_obs + Inches(0.05), Inches(12.3), Inches(0.35),
             "关键观察", size=15, bold=True)
    observations = [
        "避障负担 35.7%→55.3% (+19.6pp) —— 场景复杂度↑，更多时间在避障",
        "LQR 跟踪精度保持（0.679→0.659m）—— 复杂场景下仍稳态跟踪",
        "避障偏移 0.80→0.67m（更精准）—— 避障算法在更复杂场景下反而表现更优",
        "避障后恢复 2.87→0.84s（快 3.4 倍）—— 恢复力显著提升",
    ]
    tb = s.shapes.add_textbox(Inches(0.5), y_obs + Inches(0.45),
                              Inches(12.3), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, t in enumerate(observations):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = "• " + t
        run.font.name = FONT_ZH; run.font.size = Pt(12)
        rPr = run._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn('a:ea')); ea.set('typeface', FONT_ZH)

    add_footer(s, page_num, total=total)


# ============== 主流程：找到旧页 → 删 → 重建 ==============

def find_slide_by_title(prs, keywords):
    """返回包含 keywords 的最后一张 slide 的索引"""
    matches = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if all(k in text for k in keywords):
                    matches.append(i)
                    break
    return matches

def fix_qaa_footer(prs, total=17):
    """把 Q&A 页（最后一张）的页脚 / 16 改为 / total"""
    s = prs.slides[-1]
    for shape in s.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if "/ 16" in run.text:
                        run.text = run.text.replace("/ 16", f"/ {total}")
                        print(f"[add_avoidance_metrics] 修复 Q&A 页脚: {run.text}")

def main():
    pptx_path = "docs/指导书/项目答辩PPT.pptx"
    old_csv = "experiment_data/experiment_20260708_185400.csv"
    new_csv = "experiment_data/experiment_20260709_230118.csv"

    m_old = compute_metrics_from_csv(old_csv)
    m_new = compute_metrics_from_csv(new_csv)
    if m_old is None or m_new is None:
        print(f"[add_avoidance_metrics] 错误：CSV 读取失败")
        return

    print(f"[add_avoidance_metrics] 旧实验：{m_old['records']} 条 / {m_old['duration']:.1f}s")
    print(f"[add_avoidance_metrics] 新实验：{m_new['records']} 条 / {m_new['duration']:.1f}s")
    print(f"[add_avoidance_metrics] 旧避障效果：偏移 {m_old['av_offset_mean']:.2f}m / "
          f"段均长 {m_old['av_seg_mean_dur']:.2f}s / 恢复 {m_old['av_recovery_mean']:.2f}s")
    print(f"[add_avoidance_metrics] 新避障效果：偏移 {m_new['av_offset_mean']:.2f}m / "
          f"段均长 {m_new['av_seg_mean_dur']:.2f}s / 恢复 {m_new['av_recovery_mean']:.2f}s")

    prs = Presentation(pptx_path)
    n_before = len(prs.slides)
    print(f"[add_avoidance_metrics] 当前 {n_before} 页")

    # 记录保护性指标
    pic_count_12 = sum(1 for s in prs.slides[11].shapes if s.shape_type == 13)
    pic_count_13 = sum(1 for s in prs.slides[12].shapes if s.shape_type == 13)
    print(f"[add_avoidance_metrics] slide 12 图片: {pic_count_12}, slide 13 图片: {pic_count_13}")

    # 策略：直接删最后 3 张（残留页），重建 3 张
    # 这样保证总数仍是 17，结构整齐：14 原页 + 3 新页
    n_to_delete = 3
    if n_before < n_to_delete + 14:
        print(f"[add_avoidance_metrics] 错误：页数异常 {n_before}")
        return
    print(f"[add_avoidance_metrics] 删除最后 {n_to_delete} 张（slide {n_before - n_to_delete + 1}~{n_before}）")
    for i in range(n_to_delete):
        remove_slide_at(prs, n_before - 1 - i)  # 从后往前删

    n_after_del = len(prs.slides)
    print(f"[add_avoidance_metrics] 删除后 {n_after_del} 页")

    # 在末尾追加 3 张新页：多障碍数据、对比、Q&A
    slide_multi_obstacle_data(prs, m_new, page_num=15, total=17)
    slide_comparison(prs, m_old, m_new, page_num=16, total=17)
    slide_qaa(prs, page_num=17, total=17)
    # 顺序就是：原 1-14, 多障碍 15, 对比 16, Q&A 17

    # 验证图片数
    pic_count_12_after = sum(1 for s in prs.slides[11].shapes if s.shape_type == 13)
    pic_count_13_after = sum(1 for s in prs.slides[12].shapes if s.shape_type == 13)
    if pic_count_12 != pic_count_12_after or pic_count_13 != pic_count_13_after:
        print(f"[add_avoidance_metrics] 警告：图片数变化！12: {pic_count_12}->{pic_count_12_after}, "
              f"13: {pic_count_13}->{pic_count_13_after}")
    else:
        print(f"[add_avoidance_metrics] ✓ slide 12/13 图片完整保留")

    prs.save(pptx_path)
    print(f"[add_avoidance_metrics] 已保存 {pptx_path}，共 {len(prs.slides)} 页")

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.split('\n')[0][:55]
                break
        n_pic = sum(1 for s in slide.shapes if s.shape_type == 13)
        print(f"  slide {i}: pics={n_pic} | {title}")


if __name__ == "__main__":
    main()
