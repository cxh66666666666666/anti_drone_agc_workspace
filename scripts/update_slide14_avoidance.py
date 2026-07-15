#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
@file: update_slide14_avoidance.py
@brief: 在 slide 14 上添加避障效果 3 指标
@strategy:
    1. 清理掉所有旧的避障块（line/矩形/textbox，含多次残留）
    2. 重新紧凑排列动作分布柱状图，让出 y=4.0-5.2 给避障块
    3. 在 y=3.95-5.20 区域加：横线 + 小标题 + 3 卡片
    4. 动作分布下移到 y=5.30-6.80（标题 + 6 行 0.25 间隔）
    5. 幂等：多次运行结果一致
"""

import os
import csv
import statistics
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============== 颜色 / 字体 ==============
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY_DARK = RGBColor(0x66, 0x66, 0x66)
FONT_ZH = "Microsoft YaHei"

# ============== 避障效果计算 ==============

def compute_av_metrics(path):
    if not os.path.exists(path):
        return None
    with open(path) as f:
        rows = list(csv.DictReader(f))
    if not rows:
        return None
    is_av_b = [r["is_avoiding"] == "1" for r in rows]
    off = []
    for r in rows:
        v = r.get("offset_magnitude_xy", "nan")
        off.append(float(v) if v not in ("nan", "", None) else None)
    off_av = [off[i] for i in range(len(rows)) if is_av_b[i] and off[i] is not None]
    av_offset_mean = statistics.mean(off_av) if off_av else 0.0

    segs = []
    cur = 0
    for x in is_av_b:
        if x:
            cur += 1
        elif cur > 0:
            segs.append(cur)
            cur = 0
    if cur > 0:
        segs.append(cur)
    seg_count = len(segs)
    seg_mean_dur = (statistics.mean(segs) * 0.1) if segs else 0.0

    pes = [float(r["position_error"]) if r["position_error"] not in ("nan", "") else None
           for r in rows]
    rec = []
    for i in range(len(is_av_b) - 1):
        if is_av_b[i] and not is_av_b[i + 1]:
            for j in range(i + 1, min(i + 100, len(rows))):
                if pes[j] is not None and pes[j] < 1.0:
                    rec.append((j - i) * 0.1)
                    break
    recovery_mean = statistics.mean(rec) if rec else 0.0
    return {
        "av_offset_mean": av_offset_mean,
        "av_seg_count": seg_count,
        "av_seg_mean_dur": seg_mean_dur,
        "av_recovery_mean": recovery_mean,
    }


# ============== 绘制函数 ==============

def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_ZH
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_ZH)
    return tb


def add_rect(slide, left, top, width, height, line=BLACK, line_weight=1.0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.line.color.rgb = line
    rect.line.width = Pt(line_weight)
    rect.fill.background()
    rect.shadow.inherit = False
    rect.text_frame.text = ""
    return rect


def add_hline(slide, y, x_start=Inches(0.5), x_end=Inches(12.833), weight=1.0):
    line = slide.shapes.add_connector(1, x_start, y, x_end, y)
    line.line.color.rgb = BLACK
    line.line.width = Pt(weight)
    return line


# ============== 清理函数 ==============

def remove_existing_avoidance(slide):
    """清理所有避障效果块（line/矩形/textbox），不影响动作分布"""
    AV_KEYWORDS = ("避障效果", "避障偏移", "避障段均", "避障后恢复",
                   "is_avoiding=1 期间", "结束→err", "单次避障",
                   "共 ", "段")
    ACTION_LABELS = {"前进", "左移", "右移", "上升", "下降", "悬停"}
    to_remove = []
    for shape in list(slide.shapes):
        if not shape.top:
            continue
        top_in = shape.top / 914400
        # 横线 (y=3.85-4.10)  - LINE 类型 = 1 或 9 (connector)
        st = str(shape.shape_type)
        if (("LINE" in st) or shape.shape_type == 1) and not shape.has_text_frame:
            if 3.85 < top_in < 4.10:
                to_remove.append(shape)
            continue
        # 避障卡矩形 (4.0 width, 1.0-1.2 height) 在 y=4.0-5.5 范围
        if shape.width and shape.height and not shape.has_text_frame:
            w_in = shape.width / 914400
            h_in = shape.height / 914400
            if (abs(w_in - 4.0) < 0.5
                    and 0.8 < h_in < 1.3
                    and 3.5 < top_in < 6.0):
                to_remove.append(shape)
            continue
        # textbox
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if any(k in t for k in AV_KEYWORDS):
                to_remove.append(shape)
                continue
            # 避障卡内的纯数值/单位/说明：y in (4.0, 5.6), x in 避障卡位置
            if shape.left and 4.0 < top_in < 5.6:
                left_in = shape.left / 914400
                if (abs(left_in - 0.55) < 0.2
                        or abs(left_in - 4.75) < 0.2
                        or abs(left_in - 8.95) < 0.2):
                    if t in ACTION_LABELS or t.endswith("%"):
                        continue
                    to_remove.append(shape)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)
    return len(to_remove)


# ============== 动作分布重排 ==============

def reposition_action_distribution(slide):
    """为每个 shape 独立确定目标 action：标签按 text，百分比按值，柱子按最近邻"""
    TARGET_Y = {
        "动作分布柱状图（按当前数据集中各动作占比）": 5.30,
        "前进": 5.55,
        "左移": 5.80,
        "右移": 6.05,
        "上升": 6.30,
        "下降": 6.55,
        "悬停": 6.80,
    }
    PCT_TO_ACT = {64.3: "前进", 18.8: "左移", 16.9: "上升"}
    # shape_action 用 id(shape) 作 key
    shape_action = {}  # id(shape) -> action (str or ("BAR", y))
    shape_objs = {}   # id(shape) -> shape
    zero_pct_y_to_shape = {}  # y -> shape id（0% 行的百分比）
    for shape in slide.shapes:
        if not shape.top or not shape.left:
            continue
        top_in = shape.top / 914400
        left_in = shape.left / 914400
        if not (4.0 < top_in < 7.2):
            continue
        if not (0.4 < left_in < 11.6):
            continue
        sid = id(shape)
        shape_objs[sid] = shape
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t in TARGET_Y:
                shape_action[sid] = t
                continue
            if t.endswith('%'):
                try:
                    pct = float(t.replace('%', '').strip())
                except ValueError:
                    continue
                if pct in PCT_TO_ACT:
                    shape_action[sid] = PCT_TO_ACT[pct]
                elif pct == 0.0:
                    shape_action[sid] = None  # 待推断
                    zero_pct_y_to_shape[top_in] = sid
                continue
        # 柱子 (x=1.7) 暂时也标记为 None
        if 1.5 < left_in < 1.9 and not shape.has_text_frame:
            shape_action[sid] = ("BAR", top_in)
    # 为 0% 按 y 排序分配 action
    zero_pct_ys = sorted(zero_pct_y_to_shape.keys())
    zero_pct_action_seq = ["右移", "下降", "悬停"]
    zero_pct_y_action = {y: zero_pct_action_seq[i]
                         for i, y in enumerate(zero_pct_ys)}
    for y, sid in zero_pct_y_to_shape.items():
        shape_action[sid] = zero_pct_y_action.get(y, "右移")
    # 为柱子分配 action（找最近的已分配 action）
    assigned = []  # list of (y, action)
    for sid, act in shape_action.items():
        if isinstance(act, str) and act in TARGET_Y:
            assigned.append((shape_objs[sid].top / 914400, act))
    for sid, act in list(shape_action.items()):
        if isinstance(act, tuple) and act[0] == "BAR":
            y = act[1]
            best = None
            best_diff = 0.1
            for ay, aact in assigned:
                if abs(y - ay) < best_diff:
                    best_diff = abs(y - ay)
                    best = aact
            if best:
                shape_action[sid] = best
            else:
                del shape_action[sid]
    # 应用移动
    moved = 0
    for sid, act in shape_action.items():
        if isinstance(act, str) and act in TARGET_Y:
            shape_objs[sid].top = Inches(TARGET_Y[act])
            moved += 1
    return moved


# ============== 主流程 ==============

def main():
    pptx_path = "docs/指导书/项目答辩PPT.pptx"
    csv_path = "experiment_data/experiment_20260708_185400.csv"
    m = compute_av_metrics(csv_path)
    if m is None:
        print("[update_slide14_avoidance] CSV 读取失败")
        return
    print(f"[update_slide14_avoidance] 避障效果：偏移 {m['av_offset_mean']:.2f}m / "
          f"段均长 {m['av_seg_mean_dur']:.2f}s / 恢复 {m['av_recovery_mean']:.2f}s")

    prs = Presentation(pptx_path)
    slide = prs.slides[13]  # 第 14 页

    n_before = len(slide.shapes)

    # 1) 清理所有旧避障块
    removed = remove_existing_avoidance(slide)
    print(f"[update_slide14_avoidance] 已清理 {removed} 个旧避障 shape")

    # 2) 重新排列动作分布（让出 y=4.0-5.2 给避障块）
    moved = reposition_action_distribution(slide)
    print(f"[update_slide14_avoidance] 动作分布已重排，移动 {moved} 个 shape")

    # 3) 在 y=3.95-5.20 加避障块
    add_hline(slide, Inches(3.95), weight=0.5)
    add_text(slide, Inches(0.5), Inches(4.00), Inches(12.3), Inches(0.25),
             "避障效果（3 项指标）", size=13, bold=True, color=BLACK)

    av_cards = [
        ("避障偏移均值", f"{m['av_offset_mean']:.2f}", "m",
         f"is_avoiding=1 期间偏移 / 共 {m['av_seg_count']} 段"),
        ("避障段均长",   f"{m['av_seg_mean_dur']:.2f}", "s",
         f"单次避障平均持续 {m['av_seg_count']} 段"),
        ("避障后恢复",   f"{m['av_recovery_mean']:.2f}", "s",
         "结束→err<1m 所需时间"),
    ]
    ac_w = Inches(4.0)
    ac_gap = Inches(0.2)
    ac_x0 = Inches(0.55)
    ac_y = Inches(4.30)
    ac_h = Inches(0.80)
    for i, (t, v, u, desc) in enumerate(av_cards):
        x = ac_x0 + (ac_w + ac_gap) * i
        add_rect(slide, x, ac_y, ac_w, ac_h, line_weight=1.5)
        add_text(slide, x, ac_y + Inches(0.02), ac_w, Inches(0.22),
                 t, size=11, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x, ac_y + Inches(0.22), ac_w, Inches(0.36),
                 v, size=20, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ac_y + Inches(0.55), ac_w, Inches(0.10),
                 u, size=9, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, x, ac_y + Inches(0.65), ac_w, Inches(0.13),
                 desc, size=8, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    n_after = len(slide.shapes)
    print(f"[update_slide14_avoidance] 已添加 1 横线 + 1 标题 + 3 卡片（4×4=12 shape）")

    # 4) 验证最大 y
    max_top = 0
    for shape in slide.shapes:
        if shape.top and shape.top > max_top:
            max_top = shape.top
    print(f"[update_slide14_avoidance] slide 14 shape: {n_before} → {n_after}")
    print(f"[update_slide14_avoidance] slide 14 最大 y = {max_top/914400:.2f} in")
    if max_top >= Inches(7.5):
        print(f"[update_slide14_avoidance] 警告：超出页面底部！")

    prs.save(pptx_path)
    print(f"[update_slide14_avoidance] 已保存 {pptx_path}，共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
