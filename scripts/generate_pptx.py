#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
@file: generate_pptx.py
@brief: 完整生成项目答辩 PPT（17 页）
@pages:
    1   封面
    2   01 项目背景
    3   02 研究意义
    4   03 国内外研究现状
    5   04 项目总体目标
    6   05 系统架构
    7   06 强化学习建模
    8   07 场景设计
    9   08 训练与验证流程
    10  09 当前系统介绍
    11  10 单障碍实验数据 + 避障效果 3 指标
    12  11 多障碍实验数据
    13  12 两次实验对比
    14  13 避障效果指标
    15  14 当前在研工作
    16  15 演示材料占位
    17  16 总结与 Q&A

@strategy:
    1. 先写到 /tmp/ 验证 OK 后再 mv 到目标路径
    2. 不重复读写 PPT，一次性生成

@spec_refs:
    - create-ppt-rl-uav
    - update-ppt-progress-page
    - expand-progress-section
    - fix-position-error-filter
"""

import os
import csv
import shutil
import statistics
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============== 颜色 ==============
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
FONT_ZH = "Microsoft YaHei"

# ============== 数据源 ==============
CSV_OLD = "experiment_data/experiment_20260708_185400.csv"  # 单障碍
CSV_NEW = "experiment_data/experiment_20260709_230118.csv"  # 多障碍
TARGET_REACH_THRESHOLD = 1.0
ACT_NAMES = {0: "悬停", 1: "前进", 2: "左移", 3: "右移", 4: "上升", 5: "下降"}


# ============== 数据计算 ==============

def compute_metrics(csv_path):
    """从 CSV 计算所有需要展示的指标"""
    if not os.path.exists(csv_path):
        return None
    with open(csv_path) as f:
        rows = list(csv.DictReader(f))
    if not rows:
        return None
    n_total = len(rows)

    # 位置误差
    pes_all = [float(r["position_error"]) for r in rows
               if r["position_error"] not in ("nan", "")]
    pe_raw_mean = statistics.mean(pes_all) if pes_all else 0.0
    pe_raw_max = max(pes_all) if pes_all else 0.0

    # 过滤 A：起飞阶段
    first_reach_idx = None
    for i, r in enumerate(rows):
        pe = r["position_error"]
        if pe not in ("nan", ""):
            try:
                if float(pe) < TARGET_REACH_THRESHOLD:
                    first_reach_idx = i
                    break
            except ValueError:
                continue
    if first_reach_idx is None:
        first_reach_idx = 0

    # 过滤后样本（排除起飞 + 避障 + NaN）
    pes_filt = []
    for i, r in enumerate(rows):
        if i < first_reach_idx:
            continue
        if r["is_avoiding"] == "1":
            continue
        pe = r["position_error"]
        if pe in ("nan", ""):
            continue
        pes_filt.append(float(pe))
    pe_filt_mean = statistics.mean(pes_filt) if pes_filt else 0.0
    pe_filt_median = statistics.median(pes_filt) if pes_filt else 0.0
    pe_filt_max = max(pes_filt) if pes_filt else 0.0
    pe_filt_n = len(pes_filt)

    # 避障激活
    is_av = [r["is_avoiding"] == "1" for r in rows]
    av_n = sum(is_av)
    av_pct = av_n / n_total * 100

    # 避障段
    segs = []
    cur = 0
    for x in is_av:
        if x:
            cur += 1
        elif cur > 0:
            segs.append(cur)
            cur = 0
    if cur > 0:
        segs.append(cur)
    seg_count = len(segs)
    seg_mean_dur = (statistics.mean(segs) * 0.1) if segs else 0.0

    # 避障偏移
    off_av = []
    for i, r in enumerate(rows):
        if not is_av[i]:
            continue
        v = r.get("offset_magnitude_xy", "nan")
        if v in ("nan", "", None):
            continue
        try:
            off_av.append(float(v))
        except ValueError:
            continue
    av_offset_mean = statistics.mean(off_av) if off_av else 0.0

    # 避障后恢复时间（is_avoiding=1 → 0 后，到 position_error < 1.0）
    rec = []
    for i in range(len(is_av) - 1):
        if is_av[i] and not is_av[i + 1]:
            for j in range(i + 1, min(i + 100, len(rows))):
                pe = rows[j]["position_error"]
                if pe not in ("nan", ""):
                    try:
                        if float(pe) < TARGET_REACH_THRESHOLD:
                            rec.append((j - i) * 0.1)
                            break
                    except ValueError:
                        continue
    recovery_mean = statistics.mean(rec) if rec else 0.0

    # 动作分布（按数字 0-5）
    act_count = Counter()
    for r in rows:
        a = r.get("action", "nan")
        if a in ("nan", ""):
            continue
        try:
            act_count[int(float(a))] += 1
        except ValueError:
            continue
    act_total = sum(act_count.values())
    act_dist = {}
    for k in sorted(act_count):
        name = ACT_NAMES.get(k, f"动作{k}")
        act_dist[name] = (act_count[k], act_count[k] / act_total * 100)

    # 运行时长
    ts = [float(r["timestamp"]) for r in rows if r.get("timestamp") not in ("nan", "")]
    duration = max(ts) - min(ts) if ts else 0.0

    return {
        "n_total": n_total,
        "n_records": act_total,
        "duration": duration,
        "av_pct": av_pct,
        "av_n": av_n,
        "pe_raw_mean": pe_raw_mean,
        "pe_raw_max": pe_raw_max,
        "pe_filt_mean": pe_filt_mean,
        "pe_filt_median": pe_filt_median,
        "pe_filt_max": pe_filt_max,
        "pe_filt_n": pe_filt_n,
        "first_reach_idx": first_reach_idx,
        "seg_count": seg_count,
        "seg_mean_dur": seg_mean_dur,
        "av_offset_mean": av_offset_mean,
        "recovery_mean": recovery_mean,
        "act_dist": act_dist,
        "act_total": act_total,
    }


# ============== 绘制基础 ==============

def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT_ZH):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    return tb


def add_rect(slide, left, top, width, height, line=BLACK, fill=None, line_weight=1.0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.line.color.rgb = line
    rect.line.width = Pt(line_weight)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    rect.shadow.inherit = False
    rect.text_frame.text = ""
    return rect


def add_hline(slide, y, x_start=Inches(0.5), x_end=Inches(12.833), weight=1.0):
    line = slide.shapes.add_connector(1, x_start, y, x_end, y)
    line.line.color.rgb = BLACK
    line.line.width = Pt(weight)
    return line


def add_footer(slide, page_num, total=17, footer_left="第 12 组 · 基于强化学习的无人机自主避障与路径规划系统"):
    """页脚：左下组别小字 + 右下页码"""
    add_text(slide, Inches(0.4), Inches(7.10), Inches(8.0), Inches(0.30),
             footer_left, size=9, color=GRAY_DARK)
    add_text(slide, Inches(11.80), Inches(7.10), Inches(1.20), Inches(0.30),
             f"{page_num} / {total}", size=9, color=GRAY_DARK, align=PP_ALIGN.RIGHT)


def add_title(slide, title, sub=None):
    """页眉标题 + 副标题 + 横线"""
    add_text(slide, Inches(0.5), Inches(0.30), Inches(12.3), Inches(0.44),
             title, size=24, bold=True, color=BLACK)
    if sub:
        add_text(slide, Inches(0.5), Inches(0.90), Inches(12.3), Inches(0.35),
                 sub, size=12, color=GRAY_DARK)
    add_hline(slide, Inches(1.30), weight=0.5)


# ============== 卡片 ==============

def add_stat_card(slide, x, y, w, h, title, value, unit, desc):
    """统计卡片：边框 + 标题 + 数值 + 单位 + 说明"""
    add_rect(slide, x, y, w, h, line_weight=1.0)
    add_text(slide, x, y + Inches(0.10), w, Inches(0.30),
             title, size=12, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.40), w, Inches(0.70),
             value, size=28, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(1.10), w, Inches(0.25),
             unit, size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.35), w, Inches(0.25),
             desc, size=9, color=GRAY_DARK, align=PP_ALIGN.CENTER)


# ============== 各页内容 ==============

def slide_1_cover(prs):
    """封面"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.0),
             "基于强化学习的无人机", size=40, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.0),
             "自主避障与路径规划系统", size=40, bold=True, align=PP_ALIGN.CENTER)
    add_hline(s, Inches(4.0), weight=1.5)
    add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
             "《人工智能实践》课程项目答辩", size=20, color=GRAY_DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
             "第 12 组  ·  智科 1 班", size=14, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.4),
             "组长：陈鑫豪    组员：万前进 / 何俊峰 / 李英楷",
             size=12, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             "指导教师：叶景贞", size=12, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
             "2026 年 6 月 30 日", size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER)


def slide_2_bg(prs):
    """01 项目背景"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "01  项目背景", "无人机自主飞行在城市配送、巡检、搜救等场景需求")
    items = [
        "• 城市配送：最后 1 公里高频次、密集障碍",
        "• 电力 / 桥梁巡检：穿越复杂结构、强风扰动",
        "• 灾害搜救：GPS 拒止环境、动态未知障碍",
        "• 传统方法局限：",
        "    – 人工势场法：易陷入局部极小值",
        "    – A* / RRT：规划与控制割离，难实时",
    ]
    y = 1.8
    for it in items:
        size = 18 if it.startswith("•") else 16
        add_text(s, Inches(0.7), Inches(y), Inches(12.0), Inches(0.4),
                 it, size=size)
        y += 0.45
    add_text(s, Inches(0.7), Inches(5.2), Inches(12.0), Inches(0.5),
             "→ 引入 DQN / PPO 强化学习做端到端决策", size=20, bold=True)
    add_footer(s, 2)


def slide_3_meaning(prs):
    """02 研究意义"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "02  研究意义", "应用、技术、实践三层价值")
    cards = [
        ("应用", "城市物流 / 巡检 / 搜救", "替代人工，降低成本，提升安全"),
        ("技术", "RL + 自主避障", "填补动态未知环境下端到端决策的研究空白"),
        ("实践", "可运行仿真系统", "为后续硬件迁移与实飞验证提供基础"),
    ]
    cw = Inches(4.0)
    cx0 = Inches(0.55)
    cy = Inches(2.0)
    ch = Inches(3.5)
    for i, (t, sub, desc) in enumerate(cards):
        x = cx0 + (cw + Inches(0.2)) * i
        add_rect(s, x, cy, cw, ch, line_weight=1.5)
        add_text(s, x, cy + Inches(0.3), cw, Inches(0.5),
                 t, size=24, bold=True, align=PP_ALIGN.CENTER)
        add_hline(s, cy + Inches(0.95), x + Inches(0.5), x + cw - Inches(0.5))
        add_text(s, x + Inches(0.3), cy + Inches(1.3), cw - Inches(0.6), Inches(0.5),
                 sub, size=16, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), cy + Inches(2.1), cw - Inches(0.6), Inches(1.2),
                 desc, size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 3)


def slide_4_survey(prs):
    """03 国内外研究现状"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "03  国内外研究现状", "从传统方法到深度强化学习")
    items = [
        ("传统方法", "人工势场、A*、RRT", "结构化环境可用；动态未知环境失效"),
        ("机器学习", "SVM、决策森林", "需手工特征，泛化能力有限"),
        ("深度学习", "CNN 视觉避障", "感知端提升；决策仍依赖规则"),
        ("近年新技术", "分层 RL / MARL / Sim2Real", "复杂场景的研究热点"),
        ("前沿方向", "Transformer / MPC 融合", "长期规划与控制联合优化"),
    ]
    y = 1.6
    for t, sub, desc in items:
        add_text(s, Inches(0.7), Inches(y), Inches(2.5), Inches(0.4),
                 t, size=18, bold=True)
        add_text(s, Inches(3.5), Inches(y), Inches(4.0), Inches(0.4),
                 sub, size=14)
        add_text(s, Inches(7.7), Inches(y), Inches(5.5), Inches(0.4),
                 desc, size=12, color=GRAY_DARK)
        y += 0.7
    add_footer(s, 4)


def slide_5_goal(prs):
    """04 项目总体目标"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "04  项目总体目标", "4 模块协同：感知 + 决策 + 规划 + 控制")
    mods = [
        ("感知", "LiDAR → 12 扇区特征 + 目标向量"),
        ("决策", "DQN / PPO 选 6 离散动作"),
        ("规划", "A* / RRT 基线对比"),
        ("跟踪", "LQR 控制至目标点"),
    ]
    cw = Inches(2.85)
    cx0 = Inches(0.55)
    cy = Inches(2.5)
    ch = Inches(2.5)
    for i, (t, d) in enumerate(mods):
        x = cx0 + (cw + Inches(0.2)) * i
        add_rect(s, x, cy, cw, ch, line_weight=1.5)
        add_text(s, x, cy + Inches(0.4), cw, Inches(0.6),
                 t, size=22, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(1.2), cw, Inches(1.0),
                 d, size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
             "输入：LiDAR 点云 + 位姿 + 目标点", size=14, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
             "输出：6 维动作 + 局部目标点 → MAVROS → PX4", size=14, align=PP_ALIGN.CENTER)
    add_footer(s, 5)


def slide_6_arch(prs):
    """05 系统架构"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "05  系统架构", "ROS 多节点分层结构")
    layers = [
        ("感知层", "perception_node", "12 扇区特征"),
        ("决策层", "rl_planner_node", "启发式 / DQN"),
        ("跟踪层", "lqr_tracker_node", "LQR 控制律"),
        ("执行层", "PX4 + MAVROS", "Offboard 控制"),
    ]
    cy = Inches(2.0)
    ch = Inches(0.9)
    for i, (t, n, d) in enumerate(layers):
        y = cy + (ch + Inches(0.3)) * i
        add_rect(s, Inches(2.0), y, Inches(9.3), ch, line_weight=1.0)
        add_text(s, Inches(2.2), y + Inches(0.05), Inches(2.5), Inches(0.4),
                 t, size=16, bold=True)
        add_text(s, Inches(4.7), y + Inches(0.05), Inches(3.5), Inches(0.4),
                 n, size=14, color=GRAY_DARK)
        add_text(s, Inches(8.2), y + Inches(0.05), Inches(3.0), Inches(0.4),
                 d, size=12, color=GRAY_DARK)
        if i < len(layers) - 1:
            add_text(s, Inches(6.5), y + ch, Inches(0.3), Inches(0.3),
                     "↓", size=18, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "通信链：点云 → 扇区特征 → 动作 + 局部目标 → setpoint → PX4",
             size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 6)


def slide_7_rl(prs):
    """06 强化学习建模"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "06  强化学习建模", "状态 / 动作 / 奖励三要素")
    add_text(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(0.4),
             "状态空间（20 维）：", size=16, bold=True)
    add_text(s, Inches(1.0), Inches(2.0), Inches(12.0), Inches(0.4),
             "• 12 维扇区距离  +  2 维目标方向  +  2 维速度  +  4 维上一步动作 one-hot",
             size=14)
    add_text(s, Inches(0.7), Inches(2.6), Inches(12.0), Inches(0.4),
             "动作空间（6 维离散）：", size=16, bold=True)
    add_text(s, Inches(1.0), Inches(3.0), Inches(12.0), Inches(0.4),
             "• 悬停(0) / 前进(1) / 左移(2) / 右移(3) / 上升(4) / 下降(5)",
             size=14)
    add_text(s, Inches(0.7), Inches(3.6), Inches(12.0), Inches(0.4),
             "奖励函数：", size=16, bold=True)
    rewards = [
        ("progress", "+0.5", "接近目标"),
        ("success", "+10", "到达目标点"),
        ("collision", "-10", "碰撞障碍"),
        ("near", "-0.1", "进入近距阈值"),
        ("step", "-0.01", "时间惩罚"),
        ("smooth", "-0.05", "动作切换频繁"),
    ]
    cw = Inches(1.9)
    cx0 = Inches(0.7)
    cy = Inches(4.1)
    ch = Inches(0.9)
    for i, (t, v, d) in enumerate(rewards):
        x = cx0 + (cw + Inches(0.1)) * i
        add_rect(s, x, cy, cw, ch, line_weight=1.0)
        add_text(s, x, cy + Inches(0.05), cw, Inches(0.3),
                 t, size=11, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.32), cw, Inches(0.3),
                 v, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.60), cw, Inches(0.3),
                 d, size=8, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(0.4),
             "算法：DQN（主）/ PPO（待实现）", size=14)
    add_footer(s, 7)


def slide_8_scene(prs):
    """07 场景设计"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "07  场景设计", "4 类难度逐级提升")
    scenes = [
        ("S0", "无障碍", "基础直线飞行"),
        ("S1", "稀疏静态", "1~2 个 box 障碍"),
        ("S2", "密集静态", "5+ box + 墙"),
        ("S3", "动态障碍", "运动 box + 行人模型"),
    ]
    cw = Inches(2.85)
    cx0 = Inches(0.55)
    cy = Inches(2.0)
    ch = Inches(2.2)
    for i, (code, t, d) in enumerate(scenes):
        x = cx0 + (cw + Inches(0.2)) * i
        add_rect(s, x, cy, cw, ch, line_weight=1.5)
        add_text(s, x, cy + Inches(0.3), cw, Inches(0.5),
                 code, size=28, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.95), cw, Inches(0.4),
                 t, size=16, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(1.4), cw, Inches(0.7),
                 d, size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4),
             "YAML 配置示例：", size=14, bold=True)
    add_rect(s, Inches(1.0), Inches(5.1), Inches(11.0), Inches(1.5), line_weight=0.5)
    add_text(s, Inches(1.2), Inches(5.2), Inches(10.6), Inches(1.3),
             "scene:\n  type: S2\n  obstacles:\n    - {type: box, x: 5,  y: 0,  z: 2, size: 1.5}\n    - {type: box, x: 0,  y: 5,  z: 2, size: 1.5}\n    - {type: box, x: -5, y: -5, z: 2, size: 1.5}",
             size=11, font="Courier New")
    add_footer(s, 8)


def slide_9_train(prs):
    """08 训练与验证流程"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "08  训练与验证流程", "分阶段 + 多种子验证")
    stages = [
        ("Stage 0", "S0 直飞，验证 DQN 收敛"),
        ("Stage 1", "S1 单障碍，验证避障基本能力"),
        ("Stage 2", "S2 多障碍，验证规划能力"),
        ("Stage 3", "S3 动态障碍，验证泛化"),
        ("Stage 4", "5 颗种子 × 多回合，统计指标"),
    ]
    cy = Inches(1.7)
    ch = Inches(0.55)
    for i, (t, d) in enumerate(stages):
        y = cy + (ch + Inches(0.1)) * i
        add_rect(s, Inches(1.0), y, Inches(2.0), ch, line_weight=1.0)
        add_text(s, Inches(1.0), y, Inches(2.0), ch,
                 t, size=14, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.2), y, Inches(9.0), ch,
                 d, size=13, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             "DQN 伪代码要点：", size=14, bold=True)
    add_rect(s, Inches(1.0), Inches(5.5), Inches(11.0), Inches(1.2), line_weight=0.5)
    add_text(s, Inches(1.2), Inches(5.6), Inches(10.6), Inches(1.1),
             "for episode in range(N):\n    s = env.reset()\n    while not done:\n        a = ε-greedy(Q, s)\n        s', r, done = env.step(a)\n        replay_buffer.push(s, a, r, s', done)\n        Q ← optimize(Q, replay_buffer.sample(batch))",
             size=11, font="Courier New")
    add_footer(s, 9)


def slide_10_current_system(prs):
    """09 当前系统介绍"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "09  项目当前进展与在研工作 · 当前系统介绍",
              "Gazebo + PX4 SITL 上的 4 节点仿真系统")
    # 左栏
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.4),
             "系统组件", size=16, bold=True)
    add_hline(s, Inches(2.1), Inches(0.7), Inches(6.2))
    components = [
        "仿真平台：Gazebo + PX4 SITL + MAVROS",
        "跟踪层：LQR（已实现 8 字形轨迹可跟踪）",
        "感知层：perception_node（12 维扇区特征）",
        "规划层：rl_planner_node（启发式 6 动作）",
        "日志层：experiment_logger_node（10 Hz 写 CSV）",
        "规则式避障已能稳定工作",
    ]
    y = 2.3
    for c in components:
        add_text(s, Inches(0.7), Inches(y), Inches(5.5), Inches(0.4),
                 "• " + c, size=12)
        y += 0.4
    # 右栏
    add_text(s, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
             "ROS 通信链", size=16, bold=True)
    add_hline(s, Inches(2.1), Inches(7.0), Inches(12.5))
    ros = [
        "/velodyne_points  （激光雷达点云）",
        "/mavros/local_position/pose  （位姿）",
        "       ↓",
        "perception_node → /obstacle_features",
        "       ↓",
        "rl_planner_node → /rl_planner/local_goal",
        "       ↓",
        "lqr_tracker_node → /mavros/setpoint_position/local",
        "       ↓",
        "PX4 Offboard 控制 → Gazebo",
    ]
    y = 2.3
    for line in ros:
        size = 11 if line.startswith(" ") or "↓" in line else 12
        bold = "↓" in line
        add_text(s, Inches(7.0), Inches(y), Inches(5.5), Inches(0.35),
                 line, size=size, bold=bold)
        y += 0.35
    add_footer(s, 10)


def slide_11_single_obs(prs, m_old):
    """10 单障碍实验数据 + 避障效果 3 指标"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "10  项目当前进展与在研工作 · S1 单障碍实验数据",
              f"数据源：{CSV_OLD} · 实时统计（已排除起飞 + 避障）")
    if m_old is None:
        add_text(s, Inches(2.0), Inches(3.0), Inches(8.0), Inches(0.5),
                 "CSV 读取失败", size=20, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_footer(s, 11)
        return

    # 5 卡片
    cw = Inches(2.4)
    cx0 = Inches(0.55)
    cy = Inches(1.6)
    ch = Inches(1.7)
    cards = [
        ("运行时长", f"{m_old['duration']:.1f}", "s", "本次飞行总时长"),
        ("记录条数", f"{m_old['n_records']}", "条", "10 Hz 日志行数"),
        ("避障激活", f"{m_old['av_pct']:.1f}", "%", f"is_avoiding=1 共 {m_old['av_n']} 条"),
        ("位置误差", f"{m_old['pe_filt_mean']:.3f}", "m", f"稳态非避障期 / 样本 {m_old['pe_filt_n']}"),
        ("前进次数", f"{m_old['act_dist'].get('前进', (0,0))[0]}", "次",
         f"最常执行动作（占 {m_old['act_dist'].get('前进', (0,0))[1]:.1f}%）"),
    ]
    for i, (t, v, u, d) in enumerate(cards):
        x = cx0 + (cw + Inches(0.05)) * i
        add_stat_card(s, x, cy, cw, ch, t, v, u, d)
    # 位置误差详情
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4),
             f"位置误差详情（稳态期）："
             f"均值 {m_old['pe_filt_mean']:.3f} m  ·  "
             f"中位 {m_old['pe_filt_median']:.3f} m  ·  "
             f"最大 {m_old['pe_filt_max']:.3f} m  ·  "
             f"样本 {m_old['pe_filt_n']} 条",
             size=12)
    # 避障效果 3 指标
    add_hline(s, Inches(4.05), weight=0.5)
    add_text(s, Inches(0.5), Inches(4.10), Inches(12.3), Inches(0.3),
             "避障效果（3 项指标）", size=13, bold=True)
    av_cards = [
        ("避障偏移均值", f"{m_old['av_offset_mean']:.2f}", "m",
         f"is_avoiding=1 期间偏移 / 共 {m_old['seg_count']} 段"),
        ("避障段均长",   f"{m_old['seg_mean_dur']:.2f}", "s",
         f"单次避障平均持续 {m_old['seg_count']} 段"),
        ("避障后恢复",   f"{m_old['recovery_mean']:.2f}", "s",
         "结束→err<1m 所需时间"),
    ]
    cw2 = Inches(4.0)
    cx0_2 = Inches(0.55)
    cy2 = Inches(4.40)
    ch2 = Inches(0.95)
    for i, (t, v, u, d) in enumerate(av_cards):
        x = cx0_2 + (cw2 + Inches(0.2)) * i
        add_rect(s, x, cy2, cw2, ch2, line_weight=1.5)
        add_text(s, x, cy2 + Inches(0.05), cw2, Inches(0.22),
                 t, size=11, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy2 + Inches(0.25), cw2, Inches(0.40),
                 v, size=22, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy2 + Inches(0.62), cw2, Inches(0.15),
                 u, size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, cy2 + Inches(0.76), cw2, Inches(0.18),
                 d, size=8, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # 动作分布柱状图
    add_hline(s, Inches(5.50), weight=0.5)
    add_text(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.3),
             "动作分布柱状图（按当前数据集中各动作占比）", size=12, bold=True)
    acts = ["前进", "左移", "右移", "上升", "下降", "悬停"]
    bar_y0 = Inches(5.90)
    bar_dy = Inches(0.18)
    max_w = Inches(8.0)
    for i, an in enumerate(acts):
        cnt, pct = m_old['act_dist'].get(an, (0, 0))
        y = bar_y0 + bar_dy * i
        add_text(s, Inches(0.7), y, Inches(1.0), Inches(0.20),
                 an, size=10)
        bar_w = Emu(int(max_w * pct / 100))
        if bar_w > 0:
            add_rect(s, Inches(1.7), y, bar_w, Inches(0.16),
                     fill=BLACK, line_weight=0)
        add_text(s, Inches(9.8), y, Inches(3.0), Inches(0.20),
                 f"{cnt} 次  ·  {pct:.1f} %", size=10,
                 color=GRAY_DARK, align=PP_ALIGN.RIGHT)
    add_footer(s, 11)


def slide_12_multi_obs(prs, m_new):
    """11 多障碍实验数据"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "11  项目当前进展与在研工作 · S2 多障碍实验数据",
              f"数据源：{CSV_NEW} · 实时统计")
    if m_new is None:
        add_text(s, Inches(2.0), Inches(3.0), Inches(8.0), Inches(0.5),
                 "CSV 读取失败", size=20, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_footer(s, 12)
        return
    # 5 卡片
    cw = Inches(2.4)
    cx0 = Inches(0.55)
    cy = Inches(1.6)
    ch = Inches(1.7)
    cards = [
        ("运行时长", f"{m_new['duration']:.1f}", "s", "本次飞行总时长"),
        ("记录条数", f"{m_new['n_records']}", "条", "10 Hz 日志行数"),
        ("避障激活", f"{m_new['av_pct']:.1f}", "%", f"is_avoiding=1 共 {m_new['av_n']} 条"),
        ("位置误差", f"{m_new['pe_filt_mean']:.3f}", "m", f"稳态非避障期 / 样本 {m_new['pe_filt_n']}"),
        ("悬停次数", f"{m_new['act_dist'].get('悬停', (0,0))[0]}", "次",
         f"最常执行（占 {m_new['act_dist'].get('悬停', (0,0))[1]:.1f}%）"),
    ]
    for i, (t, v, u, d) in enumerate(cards):
        x = cx0 + (cw + Inches(0.05)) * i
        add_stat_card(s, x, cy, cw, ch, t, v, u, d)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4),
             f"位置误差详情（稳态期）："
             f"均值 {m_new['pe_filt_mean']:.3f} m  ·  "
             f"中位 {m_new['pe_filt_median']:.3f} m  ·  "
             f"最大 {m_new['pe_filt_max']:.3f} m  ·  "
             f"样本 {m_new['pe_filt_n']} 条",
             size=12)
    # 避障效果
    add_hline(s, Inches(4.05), weight=0.5)
    add_text(s, Inches(0.5), Inches(4.10), Inches(12.3), Inches(0.3),
             "避障效果（3 项指标）", size=13, bold=True)
    av_cards = [
        ("避障偏移均值", f"{m_new['av_offset_mean']:.2f}", "m",
         f"is_avoiding=1 期间 / 共 {m_new['seg_count']} 段"),
        ("避障段均长",   f"{m_new['seg_mean_dur']:.2f}", "s",
         f"单次避障平均持续 {m_new['seg_count']} 段"),
        ("避障后恢复",   f"{m_new['recovery_mean']:.2f}", "s",
         "结束→err<1m 所需时间"),
    ]
    cw2 = Inches(4.0)
    cx0_2 = Inches(0.55)
    cy2 = Inches(4.40)
    ch2 = Inches(0.95)
    for i, (t, v, u, d) in enumerate(av_cards):
        x = cx0_2 + (cw2 + Inches(0.2)) * i
        add_rect(s, x, cy2, cw2, ch2, line_weight=1.5)
        add_text(s, x, cy2 + Inches(0.05), cw2, Inches(0.22),
                 t, size=11, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy2 + Inches(0.25), cw2, Inches(0.40),
                 v, size=22, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy2 + Inches(0.62), cw2, Inches(0.15),
                 u, size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, cy2 + Inches(0.76), cw2, Inches(0.18),
                 d, size=8, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    # 动作分布
    add_hline(s, Inches(5.50), weight=0.5)
    add_text(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.3),
             "动作分布柱状图（多障碍场景）", size=12, bold=True)
    acts = ["前进", "左移", "右移", "上升", "下降", "悬停"]
    bar_y0 = Inches(5.90)
    bar_dy = Inches(0.18)
    max_w = Inches(8.0)
    for i, an in enumerate(acts):
        cnt, pct = m_new['act_dist'].get(an, (0, 0))
        y = bar_y0 + bar_dy * i
        add_text(s, Inches(0.7), y, Inches(1.0), Inches(0.20),
                 an, size=10)
        bar_w = Emu(int(max_w * pct / 100))
        if bar_w > 0:
            add_rect(s, Inches(1.7), y, bar_w, Inches(0.16),
                     fill=BLACK, line_weight=0)
        add_text(s, Inches(9.8), y, Inches(3.0), Inches(0.20),
                 f"{cnt} 次  ·  {pct:.1f} %", size=10,
                 color=GRAY_DARK, align=PP_ALIGN.RIGHT)
    add_footer(s, 12)


def slide_13_compare(prs, m_old, m_new):
    """12 两次实验对比"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "12  项目当前进展与在研工作 · 两次实验对比",
              "S1 单障碍 vs S2 多障碍 · 关键指标")
    if m_old is None or m_new is None:
        add_text(s, Inches(2.0), Inches(3.0), Inches(8.0), Inches(0.5),
                 "CSV 读取失败", size=20, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_footer(s, 13)
        return
    # 对比表
    rows = [
        ("记录条数", f"{m_old['n_records']}", f"{m_new['n_records']}"),
        ("运行时长 (s)", f"{m_old['duration']:.1f}", f"{m_new['duration']:.1f}"),
        ("避障激活 (%)", f"{m_old['av_pct']:.1f}", f"{m_new['av_pct']:.1f}"),
        ("位置误差均值 (m)",
         f"{m_old['pe_filt_mean']:.3f}", f"{m_new['pe_filt_mean']:.3f}"),
        ("位置误差最大 (m)",
         f"{m_old['pe_filt_max']:.3f}", f"{m_new['pe_filt_max']:.3f}"),
        ("避障段数", f"{m_old['seg_count']}", f"{m_new['seg_count']}"),
        ("避障段均长 (s)",
         f"{m_old['seg_mean_dur']:.2f}", f"{m_new['seg_mean_dur']:.2f}"),
        ("避障偏移均值 (m)",
         f"{m_old['av_offset_mean']:.2f}", f"{m_new['av_offset_mean']:.2f}"),
        ("避障后恢复 (s)",
         f"{m_old['recovery_mean']:.2f}", f"{m_new['recovery_mean']:.2f}"),
    ]
    # 表头
    col_x = [Inches(0.5), Inches(4.5), Inches(7.5), Inches(10.5)]
    col_w = [Inches(4.0), Inches(3.0), Inches(3.0), Inches(2.5)]
    headers = ["指标", "S1 单障碍", "S2 多障碍", "趋势"]
    y = Inches(1.6)
    for i, h in enumerate(headers):
        add_text(s, col_x[i], y, col_w[i], Inches(0.4),
                 h, size=14, bold=True, align=PP_ALIGN.CENTER)
    add_hline(s, Inches(2.05), weight=1.0)
    # 数据行
    row_h = Inches(0.40)
    y0 = Inches(2.10)
    for i, (k, v1, v2) in enumerate(rows):
        yi = y0 + row_h * i
        add_text(s, col_x[0], yi, col_w[0], row_h,
                 k, size=12)
        add_text(s, col_x[1], yi, col_w[1], row_h,
                 v1, size=12, align=PP_ALIGN.CENTER)
        add_text(s, col_x[2], yi, col_w[2], row_h,
                 v2, size=12, align=PP_ALIGN.CENTER)
        # 趋势
        try:
            n1, n2 = float(v1), float(v2)
            if abs(n1) < 1e-9:
                trend = "—"
            elif n2 > n1 * 1.05:
                trend = "↑ 升"
            elif n2 < n1 * 0.95:
                trend = "↓ 降"
            else:
                trend = "≈ 平"
        except ValueError:
            trend = "—"
        add_text(s, col_x[3], yi, col_w[3], row_h,
                 trend, size=12, align=PP_ALIGN.CENTER,
                 color=GRAY_DARK)
    # 结论
    add_hline(s, Inches(5.90), weight=0.5)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
             "结论：", size=14, bold=True)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.4),
             f"• 多障碍场景下避障激活从 {m_old['av_pct']:.1f}% 升至 {m_new['av_pct']:.1f}%，"
             f"系统对密集障碍的响应显著增强",
             size=12)
    add_text(s, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.4),
             f"• 避障段数从 {m_old['seg_count']} 段增至 {m_new['seg_count']} 段，"
             f"段均长 {m_old['seg_mean_dur']:.1f}s → {m_new['seg_mean_dur']:.1f}s，"
             f"恢复时间 {m_old['recovery_mean']:.1f}s → {m_new['recovery_mean']:.1f}s",
             size=12)
    add_footer(s, 13)


def slide_14_avoidance_metrics(prs, m_old, m_new):
    """13 避障效果指标（独立页）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "13  项目当前进展与在研工作 · 避障效果指标",
              "对比 S1 单障碍 vs S2 多障碍 的避障性能")
    if m_old is None or m_new is None:
        add_text(s, Inches(2.0), Inches(3.0), Inches(8.0), Inches(0.5),
                 "CSV 读取失败", size=20, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_footer(s, 14)
        return
    # 3 指标横排
    metrics = [
        ("避障偏移均值 (m)",
         m_old['av_offset_mean'], m_new['av_offset_mean'],
         "越低越接近期望路径"),
        ("避障段均长 (s)",
         m_old['seg_mean_dur'], m_new['seg_mean_dur'],
         "单次避障持续时间"),
        ("避障后恢复 (s)",
         m_old['recovery_mean'], m_new['recovery_mean'],
         "is_avoiding=0 后到 err<1m"),
    ]
    cw = Inches(4.0)
    cx0 = Inches(0.55)
    cy = Inches(2.0)
    ch = Inches(3.0)
    for i, (t, v1, v2, desc) in enumerate(metrics):
        x = cx0 + (cw + Inches(0.2)) * i
        add_rect(s, x, cy, cw, ch, line_weight=1.5)
        add_text(s, x, cy + Inches(0.2), cw, Inches(0.4),
                 t, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_hline(s, cy + Inches(0.65), x + Inches(0.3), x + cw - Inches(0.3))
        add_text(s, x, cy + Inches(0.8), cw, Inches(0.3),
                 "S1 单障碍", size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(1.15), cw, Inches(0.5),
                 f"{v1:.2f}", size=28, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(1.7), cw, Inches(0.3),
                 "S2 多障碍", size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(2.05), cw, Inches(0.5),
                 f"{v2:.2f}", size=28, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(2.7), cw, Inches(0.3),
                 desc, size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    # 说明
    add_hline(s, Inches(5.5), weight=0.5)
    add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
             "指标设计动机：", size=14, bold=True)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.4),
             "• 避障偏移均值：衡量避障动作的轨迹偏离程度（替代位置误差）",
             size=12)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.4),
             "• 避障段均长：单次连续避障的平均持续时间，反映响应速度",
             size=12)
    add_text(s, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.4),
             "• 避障后恢复：脱离避障后回到稳态所需时间，反映控制律稳定性",
             size=12)
    add_footer(s, 14)


def slide_15_in_progress(prs):
    """14 当前在研工作"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "14  项目当前进展与在研工作 · 当前在研工作",
              "下一步重点任务")
    cards = [
        ("多难度场景指标记录",
         "覆盖 S0~S3 全场景，每场景 5 颗随机种子 × 多回合",
         "输出：成功率 / 碰撞率 / 路径长度 / 规划时间 / 平滑度"),
        ("DQN 实现（补全 dqn_agent.py）",
         "从空文件 → 完整 Q 学习 agent",
         "组件：Q 网络 + 目标网络 + replay buffer + ε-greedy\n状态 20 维 / 动作 6 维 / 奖励塑形"),
        ("训练与对比",
         "DQN vs 当前规则式基线对比，下一步加入 A* 2D 栅格基线",
         "输出：训练曲线 + 对比表"),
    ]
    cw = Inches(12.3)
    cx = Inches(0.5)
    cy0 = Inches(1.7)
    ch = Inches(1.6)
    for i, (t, d1, d2) in enumerate(cards):
        y = cy0 + (ch + Inches(0.2)) * i
        add_rect(s, cx, y, cw, ch, line_weight=1.5)
        add_text(s, cx + Inches(0.3), y + Inches(0.1), Inches(3.0), Inches(0.4),
                 t, size=16, bold=True)
        add_text(s, cx + Inches(0.3), y + Inches(0.55), Inches(11.5), Inches(0.4),
                 d1, size=12)
        add_text(s, cx + Inches(0.3), y + Inches(0.95), Inches(11.5), Inches(0.6),
                 d2, size=11, color=GRAY_DARK)
    add_footer(s, 15)


def slide_16_demo(prs):
    """15 演示材料占位"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "15  项目当前进展与在研工作 · 演示材料占位",
              "待实验录制 / 截图后插入")
    placeholders = [
        ("[图] Gazebo 仿真场景", "3.0\" × 2.0\""),
        ("[图] RViz 轨迹与障碍", "3.0\" × 2.0\""),
        ("[视频] 规则式避障飞行回放", "3.0\" × 2.0\""),
        ("[图] 实验数据可视化", "3.0\" × 2.0\""),
    ]
    pw = Inches(5.8)
    ph = Inches(2.3)
    px0 = Inches(0.5)
    py0 = Inches(1.8)
    for i, (label, size) in enumerate(placeholders):
        col = i % 2
        row = i // 2
        x = px0 + (pw + Inches(0.3)) * col
        y = py0 + (ph + Inches(0.3)) * row
        add_rect(s, x, y, pw, ph, line_weight=1.5)
        add_text(s, x, y + ph / 2 - Inches(0.2), pw, Inches(0.5),
                 label, size=16, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, y + ph / 2 + Inches(0.3), pw, Inches(0.4),
                 size, size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             "提示：占位框可在 PowerPoint 中右键替换为图片 / 视频",
             size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 16)


def slide_17_summary(prs):
    """16 总结与 Q&A"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6),
             "16  总结与 Q & A", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_hline(s, Inches(2.9), weight=1.5)
    add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.0),
             "敬请老师批评指正", size=36, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             "第 12 组  ·  智科 1 班", size=14, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
             "陈鑫豪 / 万前进 / 何俊峰 / 李英楷", size=12,
             color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.4),
             "指导教师：叶景贞", size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 17)


# ============== 主流程 ==============

def main():
    out_path = "docs/指导书/项目答辩PPT.pptx"
    tmp_path = "/tmp/项目答辩PPT.pptx"

    # 计算数据
    print("[generate_pptx] 计算 CSV 数据...")
    m_old = compute_metrics(CSV_OLD)
    m_new = compute_metrics(CSV_NEW)
    if m_old:
        print(f"  S1 单障碍: 记录 {m_old['n_records']}, "
              f"避障 {m_old['av_pct']:.1f}%, 位置误差 {m_old['pe_filt_mean']:.3f}m")
    if m_new:
        print(f"  S2 多障碍: 记录 {m_new['n_records']}, "
              f"避障 {m_new['av_pct']:.1f}%, 位置误差 {m_new['pe_filt_mean']:.3f}m")

    # 创建 PPT（16:9, 13.33" × 7.5"）
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 17 页
    slide_1_cover(prs)
    slide_2_bg(prs)
    slide_3_meaning(prs)
    slide_4_survey(prs)
    slide_5_goal(prs)
    slide_6_arch(prs)
    slide_7_rl(prs)
    slide_8_scene(prs)
    slide_9_train(prs)
    slide_10_current_system(prs)
    slide_11_single_obs(prs, m_old)
    slide_12_multi_obs(prs, m_new)
    slide_13_compare(prs, m_old, m_new)
    slide_14_avoidance_metrics(prs, m_old, m_new)
    slide_15_in_progress(prs)
    slide_16_demo(prs)
    slide_17_summary(prs)

    # 先写到 /tmp 验证
    prs.save(tmp_path)
    print(f"[generate_pptx] 已生成 {tmp_path}, 共 {len(prs.slides)} 页")

    # 验证文件能正常读取
    try:
        prs2 = Presentation(tmp_path)
        print(f"[generate_pptx] 验证读取 OK, {len(prs2.slides)} 页")
    except Exception as e:
        print(f"[generate_pptx] 验证失败: {e}")
        return

    # mv 到目标路径
    shutil.copy2(tmp_path, out_path)
    print(f"[generate_pptx] 已复制到 {out_path}")


if __name__ == "__main__":
    main()
